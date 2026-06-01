#!/usr/bin/env python3
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

TITLE_COLOR = RGBColor(20, 40, 80)
ACCENT = RGBColor(0, 102, 153)
BODY = RGBColor(40, 40, 40)

slides = [
    ("Assignment 13", ["AI for Undefined Behaviour Detection in C/C++ Programs", "Compiler Design Lab"]),
    ("Problem Statement", ["UB in C/C++ enables aggressive compiler optimizations.", "Question: can an LLM detect UB and explain it before optimization changes behavior?"]),
    ("What is UB?", ["Examples: signed overflow, null dereference, invalid shifts, OOB access, UAF.", "Compilers assume UB does not happen.", "This can remove branches and change program behavior."]),
    ("Project Objectives", ["Identify UB patterns.", "Explain why they are UB.", "Predict LLVM optimization effects.", "Compare LLM vs existing tools."]),
    ("Benchmark Suite", ["10 UB programs.", "3 safe programs for false-positive checks.", "One main UB class per benchmark."]),
    ("LLVM Experiments", ["Compile each benchmark at O0 and O2.", "Generate LLVM IR and assembly.", "Diff outputs to show optimizer effects."]),
    ("Signed Overflow Demo", ["O0 and O2 build the same source differently.", "If x + 1 > x, LLVM may assume no signed overflow.", "The branch can disappear under optimization."]),
    ("Tools Used", ["UBSan + ASan", "Clang warnings", "clang-tidy", "cppcheck", "Gemini 2.5 Flash for LLM analysis"]),
    ("Results Summary", ["UBSan: 8/10", "Clang warnings: 4/10", "clang-tidy: 7/10", "cppcheck: 8/10", "LLM detection: 10/10"]),
    ("Key Findings", ["LLM was strongest for detection on this benchmark set.", "LLM explanations were mixed on some cases.", "No single traditional tool covered all UB classes."]),
    ("Strengths and Weaknesses", ["LLM strengths: semantic reasoning, flexible explanations.", "LLM weaknesses: no formal guarantee, some inconsistent explanations.", "Static tools: precise on known patterns but incomplete."]),
    ("Hybrid Pipeline", ["Fast static filter.", "Deep static analysis.", "LLM semantic review.", "UBSan runtime validation."]),
    ("Conclusion", ["Hybrid approach is the best practical outcome.", "Traditional tools + LLM gives better coverage than any single method."]),
    ("Future Work", ["Larger benchmark set.", "More compilers and optimization levels.", "Confidence scoring and better LLM calibration."]),
    ("Q&A", ["Thank you."])
]


def add_title(slide, title):
    title_box = slide.shapes.title
    title_box.text = title
    p = title_box.text_frame.paragraphs[0]
    if p.runs:
        run = p.runs[0]
    else:
        run = p.add_run()
    run.font.size = Pt(28)
    run.font.bold = True
    run.font.color.rgb = TITLE_COLOR


def add_bullets(slide, bullets):
    left = Inches(0.8)
    top = Inches(1.6)
    width = Inches(11.8)
    height = Inches(5.3)
    tx = slide.shapes.add_textbox(left, top, width, height)
    tf = tx.text_frame
    tf.word_wrap = True
    first = True
    for bullet in bullets:
        if first:
            p = tf.paragraphs[0]
            first = False
        else:
            p = tf.add_paragraph()
        p.text = bullet
        p.level = 0
        p.font.size = Pt(22)
        p.font.color.rgb = BODY
        p.space_after = Pt(10)
        p.bullet = True

for title, bullets in slides:
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    add_title(slide, title)
    add_bullets(slide, bullets)
    shape = slide.shapes.add_shape(1, Inches(0.8), Inches(1.18), Inches(2.6), Inches(0.08))
    shape.fill.solid()
    shape.fill.fore_color.rgb = ACCENT
    shape.line.fill.background()

prs.save("presentation/slides.pptx")
print("presentation/slides.pptx created")
